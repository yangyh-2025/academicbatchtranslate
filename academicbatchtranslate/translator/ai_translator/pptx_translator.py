# SPDX-License-Identifier: MPL-2.0
import asyncio
import regex  # [使用您依赖列表中的 regex 库]
from dataclasses import dataclass
from io import BytesIO
from typing import Self, Literal, List, Dict, Any, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.text.text import _Paragraph, TextFrame

from academicbatchtranslate.agents.segments_agent import SegmentsTranslateAgentConfig, SegmentsTranslateAgent
from academicbatchtranslate.ir.document import Document
from academicbatchtranslate.translator.ai_translator.base import AiTranslatorConfig, AiTranslator


# ---------------- 辅助工具类：语言与字体智能适配 ----------------
class LanguageHelper:
    """
    专门处理 PPTX 的语言标签与字体渲染适配。
    利用 regex 库的 Unicode 属性检测脚本类型。
    """

    # 常用语言映射 (覆盖常见写法)
    _COMMON_MAP = {
        "chinese": "zh-CN", "simplified chinese": "zh-CN", "zh": "zh-CN",
        "english": "en-US", "en": "en-US",
        "japanese": "ja-JP", "ja": "ja-JP",
        "korean": "ko-KR", "ko": "ko-KR",
        "french": "fr-FR", "fr": "fr-FR",
        "german": "de-DE", "de": "de-DE",
        "spanish": "es-ES", "es": "es-ES",
        "russian": "ru-RU", "ru": "ru-RU",
        # ... 其他语言
    }

    # [关键改进] 使用 regex 库的 Unicode 属性进行精确匹配
    # \p{Han}: 汉字
    # \p{Hiragana} / \p{Katakana}: 日文假名
    # \p{Hangul}: 韩文
    # 如果包含这些字符，说明需要启用东亚字体渲染
    _CJK_PATTERN = regex.compile(r'[\p{Han}\p{Hiragana}\p{Katakana}\p{Hangul}]')

    @classmethod
    def guess_lang_tag(cls, config_lang: str, text_content: str) -> str:
        """
        根据用户配置和实际文本内容，推断最合适的 PPT XML lang 属性。
        """
        # 1. 优先尝试解析用户配置
        if config_lang:
            clean_lang = config_lang.lower().strip()
            if clean_lang in cls._COMMON_MAP:
                return cls._COMMON_MAP[clean_lang]
            # 如果看起来像 ISO 代码 (如 'fr-FR'), 直接信赖
            if regex.match(r'^[a-z]{2,3}(-[a-z0-9]+)?$', clean_lang):
                return config_lang

        # 2. [兜底策略] 基于内容的脚本检测
        # 使用 regex 检查是否包含中日韩字符
        if cls._CJK_PATTERN.search(text_content):
            # 包含 CJK 字符 -> 声明为中文，激活东亚字体槽 (a:ea)
            # 即使是日文/韩文，设为 zh-CN 在字体回退机制上通常也能正确激活 CJK 渲染逻辑
            return "zh-CN"
        else:
            # 不含 CJK -> 默认为英文，激活西文字体槽 (a:latin)
            # 这涵盖了英文、法文、德文、俄文、越南语等绝大多数非 CJK 语言
            return "en-US"


# ---------------- 配置类 ----------------
@dataclass
class PPTXTranslatorConfig(AiTranslatorConfig):
    insert_mode: Literal["replace", "append", "prepend"] = "replace"
    separator: str = "\n"


# ---------------- 主类 ----------------
class PPTXTranslator(AiTranslator):
    """
    基于 python-pptx 的 .pptx 文件翻译器 (最终增强版)。
    使用 regex 库进行高性能的脚本检测。
    """

    def __init__(self, config: PPTXTranslatorConfig):
        super().__init__(config=config)
        self.chunk_size = config.chunk_size
        self.translate_agent = None
        self.total_chunks = 0
        glossary_dict = self.glossary.glossary_dict if self.glossary else None
        if not self.skip_translate:
            # 创建进度回调函数
            def progress_callback(current: int, total: int):
                self.total_chunks = total
                if self.progress_tracker:
                    # 计算进度百分比 (30% - 90% 区间)
                    percent = 30 + int((current / total) * 60)
                    self.progress_tracker.update(
                        percent=percent,
                        message=f"正在翻译 ({current}/{total})"
                    )

            agent_config = SegmentsTranslateAgentConfig(
                custom_prompt=config.custom_prompt, to_lang=config.to_lang, base_url=config.base_url,
                api_key=config.api_key, model_id=config.model_id, temperature=config.temperature,
                top_p=config.top_p,
                thinking=config.thinking, concurrent=config.concurrent, timeout=config.timeout,
                logger=self.logger, glossary_dict=glossary_dict, retry=config.retry,
                system_proxy_enable=config.system_proxy_enable, force_json=config.force_json,
                rpm=config.rpm,
                tpm=config.tpm,
                provider=config.provider,
                extra_body=config.extra_body,
                progress_callback=progress_callback,
            )
            self.translate_agent = SegmentsTranslateAgent(agent_config)
        self.insert_mode = config.insert_mode
        self.separator = config.separator

    # ---------------- 辅助函数：视觉样式 ----------------

    def _get_visual_style_signature(self, run) -> Tuple:
        """获取 Run 的视觉样式签名"""
        r_element = run._r
        rPr = r_element.rPr

        if rPr is None:
            return ("DEFAULT",)

        def get_bool_attr(tag_name):
            node = rPr.find(qn(f'a:{tag_name}'))
            if node is None: return None
            val = node.get('val')
            return val if val is not None else '1'

        bold = get_bool_attr('b')
        italic = get_bool_attr('i')
        u_node = rPr.find(qn('a:u'))
        underline = u_node.get('val') if u_node is not None else None
        strike_node = rPr.find(qn('a:strike'))
        strike = strike_node.get('val') if strike_node is not None else None
        sz = rPr.get('sz')
        latin = rPr.find(qn('a:latin'))
        latin_face = latin.get('typeface') if latin is not None else None
        ea = rPr.find(qn('a:ea'))
        ea_face = ea.get('typeface') if ea is not None else None

        color_sig = "INHERITED"
        for tag in ['solidFill', 'gradFill', 'noFill', 'blipFill', 'pattFill']:
            fill_node = rPr.find(qn(f'a:{tag}'))
            if fill_node is not None:
                parts = [tag]
                for child in fill_node:
                    val = child.get('val') or ""
                    parts.append(f"{child.tag.split('}')[-1]}:{val}")
                color_sig = "-".join(parts)
                break

        baseline = rPr.get('baseline')
        effect_sig = []
        for tag in ['highlight', 'effectLst', 'sp3d']:
            if rPr.find(qn(f'a:{tag}')) is not None:
                effect_sig.append(tag)

        return (bold, italic, underline, strike, sz, latin_face, ea_face, baseline, color_sig,
                tuple(sorted(effect_sig)))

    def _have_same_significant_styles(self, run1, run2) -> bool:
        """检查两个 Run 是否样式一致且紧邻"""
        if run1 is None or run2 is None: return False
        if self._get_visual_style_signature(run1) != self._get_visual_style_signature(run2): return False
        try:
            r1_element = run1._r
            r2_element = run2._r
            parent = r1_element.getparent()
            if parent != r2_element.getparent(): return False
            if parent.index(r2_element) != parent.index(r1_element) + 1: return False
        except Exception:
            return False
        return True

    def _apply_lang_correction(self, run, text_content: str):
        """[智能修正] 根据配置和文本内容，设置正确的 lang 属性"""
        if not text_content: return
        best_lang = LanguageHelper.guess_lang_tag(self.config.to_lang, text_content)
        if best_lang:
            rPr = run._r.get_or_add_rPr()
            rPr.set('lang', best_lang)
            rPr.set('altLang', best_lang)

    # ---------------- 核心遍历逻辑 ----------------

    def _process_text_frame(self, text_frame: TextFrame, elements: List[Dict[str, Any]], texts: List[str]):
        for paragraph in text_frame.paragraphs:
            self._process_paragraph(paragraph, elements, texts)

    def _process_paragraph(self, paragraph: _Paragraph, elements: List[Dict[str, Any]], texts: List[str]):
        if not paragraph.runs: return

        state = {'current_runs': []}

        def flush_segment():
            current_runs = state['current_runs']
            if not current_runs: return
            full_text = "".join(r.text for r in current_runs)
            if full_text.strip():
                elements.append({
                    "type": "text_runs",
                    "runs": list(current_runs),
                    "paragraph": paragraph,
                    "text_frame": paragraph._parent
                })
                texts.append(full_text)
            current_runs.clear()

        for run in paragraph.runs:
            if not run.text: continue
            last_run = state['current_runs'][-1] if state['current_runs'] else None
            if last_run and not self._have_same_significant_styles(last_run, run):
                flush_segment()
            state['current_runs'].append(run)

        flush_segment()

    def _process_shape(self, shape, elements: List[Dict[str, Any]], texts: List[str]):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child_shape in shape.shapes:
                self._process_shape(child_shape, elements, texts)
            return

        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if hasattr(cell, "text_frame") and cell.text_frame:
                        self._process_text_frame(cell.text_frame, elements, texts)
            return

        if shape.has_text_frame:
            try:
                self._process_text_frame(shape.text_frame, elements, texts)
            except Exception:
                pass

    def _scan_deep_xml_for_text(self, slide_element, elements: List[Dict[str, Any]], texts: List[str]):
        MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
        MC_ALT = f"{{{MC_NS}}}AlternateContent"
        MC_CHOICE = f"{{{MC_NS}}}Choice"
        P_SP = qn('p:sp')
        P_TXBODY = qn('p:txBody')

        for alt_content in slide_element.iter(MC_ALT):
            choice = alt_content.find(MC_CHOICE)
            if choice is None: continue
            for sp in choice.iter(P_SP):
                txBody = sp.find(P_TXBODY)
                if txBody is not None:
                    try:
                        tf = TextFrame(txBody, None)
                        self._process_text_frame(tf, elements, texts)
                    except Exception as e:
                        self.logger.warning(f"Deep XML Scan Error: {e}")

    def _scan_presentation_content(self, prs: Presentation, elements: List[Dict[str, Any]], texts: List[str]):
        def scan_slide_object(slide_obj):
            for shape in slide_obj.shapes:
                self._process_shape(shape, elements, texts)
            self._scan_deep_xml_for_text(slide_obj.element, elements, texts)

        for slide in prs.slides:
            scan_slide_object(slide)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                self._process_text_frame(slide.notes_slide.notes_text_frame, elements, texts)

        for master in prs.slide_masters:
            scan_slide_object(master)
            for layout in master.slide_layouts:
                scan_slide_object(layout)

    # ---------------- 翻译逻辑 ----------------

    def _pre_translate(self, document: Document) -> Tuple[Presentation, List[Dict[str, Any]], List[str]]:
        prs = Presentation(BytesIO(document.content))
        elements, texts = [], []
        self._scan_presentation_content(prs, elements, texts)
        self.logger.info(f"Extracted {len(texts)} text segments.")
        return prs, elements, texts

    def _apply_translation(self, element_info: Dict[str, Any], final_text: str):
        runs = element_info["runs"]
        if not runs: return

        original_text = "".join(r.text for r in runs)
        text_to_set = final_text
        if self.insert_mode == "append":
            text_to_set = original_text + self.separator + final_text
        elif self.insert_mode == "prepend":
            text_to_set = final_text + self.separator + original_text

        primary_run = runs[0]
        try:
            primary_run.text = text_to_set
            # 调用利用 regex 的智能修正
            self._apply_lang_correction(primary_run, text_to_set)

            text_frame = element_info.get("text_frame")
            if text_frame and hasattr(text_frame, 'auto_size'):
                if text_frame.auto_size == MSO_AUTO_SIZE.NONE:
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception as e:
            self.logger.warning(f"Error applying translation: {e}")
            return

        for i in range(1, len(runs)):
            runs[i].text = ""

    def _after_translate(self, prs: Presentation, elements: List[Dict[str, Any]], translated: List[str],
                         originals: List[str]) -> bytes:
        if len(elements) != len(translated):
            min_len = min(len(elements), len(translated))
            elements = elements[:min_len]
            translated = translated[:min_len]

        for info, trans in zip(elements, translated):
            self._apply_translation(info, trans)

        output_stream = BytesIO()
        prs.save(output_stream)
        return output_stream.getvalue()

    # ---------------- 接口 ----------------

    def translate(self, document: Document) -> Self:
        prs, elements, originals = self._pre_translate(document)
        if not originals:
            self.logger.info("No text found.")
            document.content = self._after_translate(prs, elements, [], [])
            return self

        if self.glossary_agent:
            glossary_dict_gen = self.glossary_agent.send_segments(originals, self.chunk_size)
            if self.glossary: self.glossary.update(glossary_dict_gen)
            if self.translate_agent and self.glossary: self.translate_agent.update_glossary_dict(
                self.glossary.glossary_dict)

        translated = self.translate_agent.send_segments(originals,
                                                        self.chunk_size) if self.translate_agent else originals
        document.content = self._after_translate(prs, elements, translated, originals)
        return self

    async def translate_async(self, document: Document) -> Self:
        prs, elements, originals = await asyncio.to_thread(self._pre_translate, document)
        if not originals:
            self.logger.info("No text found.")
            document.content = await asyncio.to_thread(self._after_translate, prs, elements, [], [])
            return self

        if self.glossary_agent:
            glossary_dict_gen = await self.glossary_agent.send_segments_async(originals, self.chunk_size)
            if self.glossary: self.glossary.update(glossary_dict_gen)
            if self.translate_agent and self.glossary: self.translate_agent.update_glossary_dict(
                self.glossary.glossary_dict)

        translated = await self.translate_agent.send_segments_async(originals,
                                                                    self.chunk_size) if self.translate_agent else originals
        document.content = await asyncio.to_thread(self._after_translate, prs, elements, translated, originals)
        return self